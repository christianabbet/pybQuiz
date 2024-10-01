import os
from typing import Optional
from tqdm import tqdm
from pptx import Presentation     
from pptx.util import Mm
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pathlib import Path
from sys import platform
import pptx
import numpy as np
import string

from pybquiz.export.base import Export
from pybquiz.background import BackgroundManager
from pybquiz.const import Const as C
from pybquiz.const import TriviaConst as TC
from pybquiz.const import WWTBAMConst as WC

from pybquiz.export.slidetemplate import SlideTemplate


class PPTXExport(Export):
    
    # Class constants
    TITLE_SUBTITLE = 0
    TITLE_CONTENT = 1
    SECTION_HEADER = 2
    TWO_CONTENT = 3
    COMPARISON = 4
    TITLE = 5
    BLANK = 6
    CONTENT_CAPTION = 7
    IMAGE_CAPTION = 8

    # Font
    FONT_TITLE = 40
    FONT_SUBTITLE = 28
    FONT_QUESTION = 28
    FONT_ANSWER = 20

    # Font location linux
    FONT_PATH = "/usr/share/fonts"
    
    def __init__(
        self, 
        dump_path: str, 
        dirout: str,
        width: Optional[float] = 254, 
        height: Optional[float] = 190.5,
        dircache: Optional[str] = ".cache",
    ) -> None:
        
        # Call super
        super().__init__(dump_path=dump_path, dircache=dircache)
        
        # Init local
        self.dirout = dirout
        self.fileout = os.path.join(dirout, os.path.splitext(os.path.basename(dump_path))[0] + ".pptx")

        self.width = width
        self.height = height
        self.root = Presentation() 
        self.st = SlideTemplate(width=width, height=height)
        self.pfont = self._get_system_fonts()
    
    def make_title(self, title: str, subtitle: str, img_bg: str, type: Optional[str] = None):
        
        color_bg=RGBColor.from_string(self.st.TRIVIA_COLOR_BBOX_BACKGROUND)
        color_text=RGBColor.from_string(self.st.TRIVIA_COLOR_TEXT)
        color_line=RGBColor.from_string(self.st.TRIVIA_COLOR_BBOX_LINE)
        
        if type is not None and type == "wwtbam":
            color_bg=RGBColor.from_string(self.st.WWTBAM_COLOR_BBOX_BACKGROUND)
            color_text=RGBColor.from_string(self.st.WWTBAM_COLOR_TEXT)
            color_line=RGBColor.from_string(self.st.WWTBAM_COLOR_BBOX_LINE)

        self._add_title_subtitle(
            title=title,
            subtitle=subtitle,
            img_bg=img_bg,
            pfont=self.pfont,
            color_bg=color_bg,
            color_text=color_text,
            color_line=color_line,
        )
                    
    def make_question(self, data: dict, prefix: str, show_answer: bool, type: str, img_bg: str, img_bg_blur: str):

        if type == "trivia":
            self._add_trivia_question(data=data, prefix=prefix, show_answer=show_answer, img_bg=img_bg_blur)
        elif type == "wwtbam":
            self._add_wwtbam_question(data=data, prefix=prefix, show_answer=show_answer, img_bg=img_bg)
        else:
            raise NotImplementedError
        
    def save(self):
        # Saving file 
        self.root.save(self.fileout)  
        
    def _add_wwtbam_question(self, data: dict, prefix: str, show_answer: bool, img_bg: str):
        
        # Extract infos
        question = data.get(TC.KEY_QUESTION, C.KEY_ERROR)
        question_value = "{}: ₿ {:,.0f}".format(prefix, data.get(WC.KEY_VALUE, "0"))
        answers_order = data.get(TC.EXT_KEY_ORDER, [])
        answers = [data.get(a, None) for a in answers_order]
        answers_id = data.get(TC.EXT_KEY_ORDER_ID, [])

        # Add slide and question            
        question_slide = self.root.slides.add_slide(self.root.slide_layouts[PPTXExport.BLANK]) 
        
        # Add text shape
        shapes =  question_slide.shapes

        # Add background
        if img_bg is not None:
            # Get new base name
            img_base, img_ext = os.path.splitext(img_bg)
            if not show_answer:
                img_bg = "{}_q{}".format(img_base, img_ext)
            else:
                img_bg = "{}_q{}{}".format(img_base, answers_id, img_ext)
            shapes.add_picture(img_bg, Mm(0), Mm(0), height=Mm(self.height))
    
        PPTXExport._add_frame(
            shapes=shapes,
            bbox=[Mm(v) for v in self.st.WWTBAM_BBOX_QT],
            text = question,
            color_bg=None,
            color_line=None,                 
            color_text=RGBColor.from_string(self.st.WWTBAM_COLOR_TEXT),                     
            font=PPTXExport.FONT_QUESTION,
            alignement=PP_ALIGN.LEFT,
            font_file=self.pfont,
        )
        
                    
        # Question
        PPTXExport._add_frame(
            shapes=shapes,
            bbox=[Mm(v) for v in self.st.WWTBAM_BBOX_VALUE],
            text = question_value,
            color_bg=None,
            color_line=None,  
            color_text=RGBColor.from_string(self.st.WWTBAM_COLOR_TEXT),                  
            font=PPTXExport.FONT_ANSWER,
            alignement=PP_ALIGN.CENTER,
            font_file=self.pfont,
        )
        
        # Set hard coded poistions
        positions_key = [
            self.st.WWTBAM_BBOX_QA, self.st.WWTBAM_BBOX_QB, 
            self.st.WWTBAM_BBOX_QC, self.st.WWTBAM_BBOX_QD,
        ]
        
        positions_value = [
            self.st.WWTBAM_BBOX_QAV, self.st.WWTBAM_BBOX_QBV, 
            self.st.WWTBAM_BBOX_QCV, self.st.WWTBAM_BBOX_QDV,
        ]
        # iterate over letter for questions
        for i in range(4):
            
            # Letter
            PPTXExport._add_frame(
                shapes=shapes,
                bbox=[Mm(v) for v in positions_key[i]],
                text = string.ascii_uppercase[i],
                color_bg=None,
                color_line=None,  
                color_text=RGBColor.from_string(self.st.WWTBAM_COLOR_QTEXT),                  
                font=PPTXExport.FONT_QUESTION,
                alignement=PP_ALIGN.LEFT,
                font_file=self.pfont,
            )
            
            # Question
            PPTXExport._add_frame(
                shapes=shapes,
                bbox=[Mm(v) for v in positions_value[i]],
                text = answers[i],
                color_bg=None,
                color_line=None,  
                color_text=RGBColor.from_string(self.st.WWTBAM_COLOR_TEXT),                  
                font=PPTXExport.FONT_QUESTION,
                alignement=PP_ALIGN.LEFT,
                font_file=self.pfont,
            )
        
        
    
    def _add_trivia_question(self, data: dict, prefix: str, show_answer: bool, img_bg: str):
            
        # Extract infos
        question = data.get(TC.KEY_QUESTION, C.KEY_ERROR)
        difficulty = data.get(TC.KEY_DIFFICULTY, C.KEY_ERROR)
        answers_order = data.get(TC.EXT_KEY_ORDER, [])
        answers = [data.get(a, None) for a in answers_order]
        answers_id = data.get(TC.EXT_KEY_ORDER_ID, [])
        
        # Add slide and question            
        question_slide = self.root.slides.add_slide(self.root.slide_layouts[PPTXExport.BLANK]) 
        
        # Add text shape
        shapes =  question_slide.shapes

        if img_bg is not None:
            shapes.add_picture(img_bg, Mm(0), Mm(0), height=Mm(self.height))
            
        # Add question
        (x, y, w, h) = self.st.get_trivia_question_bbox()
        PPTXExport._add_frame(
            shapes=shapes,
            bbox=[Mm(x), Mm(y), Mm(w), Mm(h)],
            text=None,
            font_file=self.pfont,
            color_bg=RGBColor.from_string(self.st.TRIVIA_COLOR_BBOX_BACKGROUND), 
            color_text=RGBColor.from_string(self.st.TRIVIA_COLOR_TEXT), 
            color_line=RGBColor.from_string(self.st.TRIVIA_COLOR_BBOX_LINE), 
        )
        
        if difficulty is not None:
            # Add difficulty
            (x, y, w, h) = self.st.get_difficulty_bbox()
            PPTXExport._add_frame(
                shapes=shapes,
                bbox=[Mm(x), Mm(y), Mm(w), Mm(h)],
                text=None,
                font_file=self.pfont,
                shapeid=MSO_SHAPE.OVAL,         
                color_bg=RGBColor.from_string(self.st.COLOR_DIFFICULTY[difficulty]), 
                color_line=RGBColor.from_string(self.st.COLOR_DIFFICULTY[difficulty]), 
            )
                    
        (x, y, w, h) = self.st.get_trivia_question_bbox()
        m_in = self.st.get_inner_margin()
        PPTXExport._add_frame(
            shapes=shapes,
            bbox=[Mm(x + m_in), Mm(y + m_in), Mm(w - 2*m_in), Mm(h - 2*m_in)],
            text = "{}: {}".format(prefix, question),
            color_bg=None,
            color_line=None,                    
            font=PPTXExport.FONT_QUESTION,
            alignement=PP_ALIGN.LEFT,
            font_file=self.pfont,
        )
        
        self._add_trivia_answers(question_slide, answers, answers_id, show_answer)       
            
             
    def _get_system_fonts(self, template: str = "Amiri-Regular"):
        # On linux need to look for font
        if "linux" in platform:
            # Find fonts
            fonts = [str(r) for r in Path(self.FONT_PATH).rglob("*.ttf")]
            tfonts = [f for f in fonts if "{}.ttf".format(template) in f]
            if len(tfonts) == 0:
                # Not font, take random
                pfont = fonts[0]
            else:
                # Take tmeplate font
                pfont = tfonts[0]
            # Look for template
            return pfont
        else:
            return None
            
    def _add_trivia_answers(self, slide, answers: list[str], answers_id: int = None, show_answer: bool = False):
                           
            nQ = len(answers)
            nRows = np.ceil(nQ/2).astype(int)
                                    
            # Check if only one proposition and no answer
            if nQ <= 1 and not show_answer:
                return
            
            # Get total heigh based on the number of entries
            answer_height, answer_inter_height = self.st.get_answer_heights()
            im = self.st.get_iinner_margin()
            col, col1, col2 = self.st.get_answer_cols()
            htot = nRows * answer_height + (nRows-1) * answer_inter_height
            (_, ay, _, ah) = self.st.get_answer_bbox()
            hoffset = (ah - htot) / 2
            
            # Add propositions
            for j in range(nQ):
                # Define col
                id_col = j % 2
                id_row = j // 2
                
                # Fix width
                if id_col == 0:
                    left = col1
                else:
                    left = col2
            
                # Add background
                local_offset = id_row * (answer_height + answer_inter_height)
                shapes =  slide.shapes
                
                # Check if one or more answers
                a = answers[j]
    
                if show_answer and j == answers_id:
                    color_bg = RGBColor.from_string(self.st.TRIVIA_COLOR_BBOX_BACKGROUND_CORRECT)
                else:
                    color_bg = RGBColor.from_string(self.st.TRIVIA_COLOR_BBOX_BACKGROUND)
                    
                self._add_frame(
                    shapes=shapes,
                    bbox=[                
                        Mm(left), 
                        Mm(local_offset + ay + hoffset), 
                        Mm(col), 
                        Mm(answer_height)
                    ],
                    text = None,
                    color_bg=color_bg,
                    color_line=color_bg,
                    font_file=self.pfont,
                )
                            
                self._add_frame(
                    shapes=shapes,
                    bbox=[                
                        Mm(left + im), 
                        Mm(local_offset + ay + hoffset + im), 
                        Mm(col - 2*im), 
                        Mm(answer_height - 2*im)                        
                    ],
                    text=a,
                    color_bg=None,
                    color_line=None,
                    font=PPTXExport.FONT_ANSWER,
                    alignement=PP_ALIGN.LEFT,
                    font_file=self.pfont,
                )
                                
                                            
    def _add_title_subtitle(
        self, title: str = "", 
        subtitle: str = None, 
        img_bg: str = None, 
        pfont: str = None,
        color_bg = None,
        color_text = None,
        color_line = None,
    ):
        
        # Add slide
        title_slide = self.root.slides.add_slide(self.root.slide_layouts[PPTXExport.BLANK]) 
        shapes =  title_slide.shapes            
                    
        # Add background
        if img_bg is not None:
            shapes.add_picture(img_bg, Mm(0), Mm(0), height=Mm(self.height))
        
        (x, y, w, h) = self.st.get_title_bbox()
        PPTXExport._add_frame(
            shapes=shapes,
            bbox=[Mm(x), Mm(y), Mm(w), Mm(h)],
            text=title,
            font=PPTXExport.FONT_TITLE,
            font_file=pfont,
            color_bg=color_bg, 
            color_text=color_text, 
            color_line=color_line, 
        )
        
        if subtitle is not None:
            (x, y, w, h) = self.st.get_subtitle_bbox()
            PPTXExport._add_frame(
                shapes=shapes,
                bbox=[Mm(x), Mm(y), Mm(w), Mm(h)],
                text=subtitle,
                font= PPTXExport.FONT_SUBTITLE,
                font_file=pfont,
                color_bg=color_bg, 
                color_text=color_text, 
                color_line=color_line, 
            )
            
    @staticmethod
    def _add_frame(
        shapes: pptx.shapes.shapetree.SlideShapes, 
        bbox: list[int], 
        text: str, 
        font: int = 0, 
        color_bg: str = None, 
        color_text: str = None, 
        color_line: str = None, 
        shapeid: int = MSO_SHAPE.ROUNDED_RECTANGLE,         
        alignement: int = PP_ALIGN.CENTER,
        font_file: str = None,
    ):
        
        shape = shapes.add_shape(
            shapeid, 
            bbox[0], bbox[1], bbox[2], bbox[3],
        )
        shape.shadow.inherit = False            
        tf = shape.text_frame
        
        if color_bg is not None:
            shape.fill.solid()
            shape.fill.fore_color.rgb = color_bg
        else:
            shape.fill.background()
            
        if color_line is not None:
            shape.line.fill.solid()
            shape.line.fill.fore_color.rgb = color_line
        else:
            shape.line.fill.background()

        if text is not None:
            tf.text = text
            tf.fit_text(max_size=font, font_file=font_file) 
            tf.paragraphs[0].alignment = alignement
        
        if color_text is not None:
            tf.paragraphs[0].font.color.rgb = color_text
        