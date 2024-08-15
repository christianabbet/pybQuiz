from pptx.util import Mm
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
import pptx
from pptx import Presentation     
import numpy as np
import json
from pybquiz.const import PYBConst as C
from sys import platform
import os
from pathlib import Path
from pybquiz.export.slidetemplate import SlideTemplate


def pos_to_char(pos):
    return chr(pos + 97)


class PptxFactory:
        
    TITLE_SUBTITLE = 0
    TITLE_CONTENT = 1
    SECTION_HEADER = 2
    TWO_CONTENT = 3
    COMPARISON = 4
    TITLE = 5
    BLANK = 6
    CONTENT_CAPTION = 7
    IMAGE_CAPTION = 8

    # Font
    FONT_TITLE = 40
    FONT_SUBTITLE = 28
    FONT_QUESTION = 28
    FONT_ANSWER = 20

    # Font location linux
    FONT_PATH = "/usr/share/fonts"
    
    def __init__(self, width: float = 254, height: float = 190.5) -> None:
        # Save results
        self.width = width
        self.height = height
        self.root = Presentation() 
        self.st = SlideTemplate(width=width, height=height)
        
    def export(self, dump_path: str, outfile: str, background_gen = None):
        
        # Reload data
        with open(dump_path, "r") as f:
            data = json.load(f)

        # Get font reference
        pfont = PptxFactory.get_system_fonts()

        # Add splide to presetation
        img_bg_quiz = background_gen.get_background("Pub Quiz")
        img_bg_paper = background_gen.get_background("Quiz paper sheets")
        self._add_title_subtitle(title=data[C.TITLE], subtitle=data[C.AUTHOR], img_bg=img_bg_quiz, pfont=pfont) 
                    
        # Create rounds
        Nround = len(data[C.ROUNDS])
        for i in range(Nround):
            # Create title slide
            catname = data[C.ROUNDS][i][C.QUESTIONS][0][C.CATEGORY]
            img_bg = background_gen.get_background(catname, blurred=False)
            img_bg_blurred = background_gen.get_background(catname, blurred=True)
            str_round = "Round {}: {}".format(i+1, data[C.ROUNDS][i][C.TITLE])
            self._add_title_subtitle(title=str_round, img_bg=img_bg, pfont=pfont)            
            
            Nquestion = len(data[C.ROUNDS][i][C.QUESTIONS])
            for j in range(Nquestion):
                # Add question slide
                answers = data[C.ROUNDS][i][C.QUESTIONS][j][C.ANSWERS]
                difficulty = data[C.ROUNDS][i][C.QUESTIONS][j][C.DIFFICULTY]
                self._add_question(
                    i=j+1, 
                    question=data[C.ROUNDS][i][C.QUESTIONS][j][C.QUESTIONS], 
                    answers=answers,
                    difficulty=difficulty,
                    img_bg=img_bg_blurred,
                    pfont=pfont,
                )                

            self._add_title_subtitle(title="Exchange paper sheets", subtitle=None, img_bg=img_bg_paper, pfont=pfont)          
            self._add_title_subtitle(title=str_round, subtitle="Answers", img_bg=img_bg, pfont=pfont)            
               
            # Iterate over questions (answers)
            for j in range(Nquestion):
                # Add question slide
                answers = data[C.ROUNDS][i][C.QUESTIONS][j][C.ANSWERS]
                difficulty = data[C.ROUNDS][i][C.QUESTIONS][j][C.DIFFICULTY]                
                correct_answers = data[C.ROUNDS][i][C.QUESTIONS][j][C.CORRECT_ANSWERS]
                self._add_question(
                    i=j+1, 
                    question=data[C.ROUNDS][i][C.QUESTIONS][j][C.QUESTIONS], 
                    answers=answers,
                    difficulty=difficulty,                    
                    answers_id=correct_answers,
                    img_bg=img_bg_blurred,
                    pfont=pfont,
                )                 

            # Add splide to presetation
            self._add_title_subtitle(title="Bring back paper sheets", subtitle=None, img_bg=img_bg_paper, pfont=pfont) 
            
        # Saving file 
        self.root.save(outfile)             
        
    @staticmethod
    def get_system_fonts(template: str = "Amiri-Regular"):
        # On linux need to look for font
        if "linux" in platform:
            # Find fonts
            fonts = [str(r) for r in Path(PptxFactory.FONT_PATH).rglob("*.ttf")]
            tfonts = [f for f in fonts if "{}.ttf".format(template) in f]
            if len(tfonts) == 0:
                # Not font, take random
                pfont = fonts[0]
            else:
                # Take tmeplate font
                pfont = tfonts[0]
            # Look for template
            return pfont
        else:
            return None


    @staticmethod
    def _add_frame(
        shapes: pptx.shapes.shapetree.SlideShapes, 
        bbox: list[int], 
        text: str, 
        font: int = 0, 
        color_bg: str = None, 
        color_text: str = None, 
        color_line: str = None, 
        shapeid: int = MSO_SHAPE.ROUNDED_RECTANGLE,         
        alignement: int = PP_ALIGN.CENTER,
        font_file: str = None,
    ):
        
        shape = shapes.add_shape(
            shapeid, 
            bbox[0], bbox[1], bbox[2], bbox[3],
        )
        shape.shadow.inherit = False            
        tf = shape.text_frame
        
        if color_bg is not None:
            shape.fill.solid()
            shape.fill.fore_color.rgb = color_bg
        else:
            shape.fill.background()
            
        if color_line is not None:
            shape.line.fill.solid()
            shape.line.fill.fore_color.rgb = color_line
        else:
            shape.line.fill.background()

        if text is not None:
            tf.text = text
            tf.fit_text(max_size=font, font_file=font_file) 
            tf.paragraphs[0].alignment = alignement
        
        if color_text is not None:
            tf.paragraphs[0].font.color.rgb = color_text
        
    def _add_title_subtitle(self, title: str = "", subtitle: str = None, img_bg: str = None, pfont: str = None):
        
            # Add slide
            title_slide = self.root.slides.add_slide(self.root.slide_layouts[PptxFactory.BLANK]) 
            shapes =  title_slide.shapes            
                        
            # Add background
            if img_bg is not None:
                shapes.add_picture(img_bg, Mm(0), Mm(0), height=Mm(self.height))
            
            (x, y, w, h) = self.st.get_title_bbox()
            PptxFactory._add_frame(
                shapes=shapes,
                bbox=[Mm(x), Mm(y), Mm(w), Mm(h)],
                text=title,
                font=PptxFactory.FONT_TITLE,
                font_file=pfont,
                color_bg=RGBColor.from_string(self.st.COLOR_BBOX_BACKGROUND), 
                color_text=RGBColor.from_string(self.st.COLOR_TEXT), 
                color_line=RGBColor.from_string(self.st.COLOR_BBOX_LINE), 
            )
            
            if subtitle is not None:
                (x, y, w, h) = self.st.get_subtitle_bbox()
                PptxFactory._add_frame(
                    shapes=shapes,
                    bbox=[Mm(x), Mm(y), Mm(w), Mm(h)],
                    text=subtitle,
                    font= PptxFactory.FONT_SUBTITLE,
                    font_file=pfont,
                    color_bg=RGBColor.from_string(self.st.COLOR_BBOX_BACKGROUND), 
                    color_text=RGBColor.from_string(self.st.COLOR_TEXT), 
                    color_line=RGBColor.from_string(self.st.COLOR_BBOX_LINE), 
                )
            
    def _add_question(self, i: int = 1, question: str = "", difficulty: int = None, answers: list[str] = [], answers_id: list[int] = None, img_bg: str = None, pfont: str = None):
        
            # Add slide and question            
            question_slide = self.root.slides.add_slide(self.root.slide_layouts[PptxFactory.BLANK]) 
            
            # Add text shape
            shapes =  question_slide.shapes

            if img_bg is not None:
                shapes.add_picture(img_bg, Mm(0), Mm(0), height=Mm(self.height))
                
            # Add question
            (x, y, w, h) = self.st.get_question_bbox()
            PptxFactory._add_frame(
                shapes=shapes,
                bbox=[Mm(x), Mm(y), Mm(w), Mm(h)],
                text=None,
                font_file=pfont,
                color_bg=RGBColor.from_string(self.st.COLOR_BBOX_BACKGROUND), 
                color_text=RGBColor.from_string(self.st.COLOR_TEXT), 
                color_line=RGBColor.from_string(self.st.COLOR_BBOX_LINE), 
            )

            if difficulty is not None:
                # Add difficulty
                (x, y, w, h) = self.st.get_difficulty_bbox()
                PptxFactory._add_frame(
                    shapes=shapes,
                    bbox=[Mm(x), Mm(y), Mm(w), Mm(h)],
                    text=None,
                    font_file=pfont,
                    shapeid=MSO_SHAPE.OVAL,         
                    color_bg=RGBColor.from_string(self.st.COLOR_DIFFICULTY[difficulty]), 
                    color_line=RGBColor.from_string(self.st.COLOR_DIFFICULTY[difficulty]), 
                )
                        
            (x, y, w, h) = self.st.get_question_bbox()
            m_in = self.st.get_inner_margin()
            PptxFactory._add_frame(
                shapes=shapes,
                bbox=[Mm(x + m_in), Mm(y + m_in), Mm(w - 2*m_in), Mm(h - 2*m_in)],
                text = "Q{}: {}".format(i, question),
                color_bg=None,
                color_line=None,                    
                font=PptxFactory.FONT_QUESTION,
                alignement=PP_ALIGN.LEFT,
                font_file=pfont,
            )

            self._add_answers(question_slide, answers, answers_id, pfont=pfont)

            
    def _add_answers(self, slide, answers: list[str], answers_id: list[int] = None, pfont: str = None):
                           
            nQ = len(answers)
            nRows = np.ceil(nQ/2).astype(int)
                                    
            # Check if only one proposition and no answer
            if nQ <= 1 and answers_id is None:
                return
            
            # Get total heigh based on the number of entries
            answer_height, answer_inter_height = self.st.get_answer_heights()
            im = self.st.get_iinner_margin()
            col, col1, col2 = self.st.get_answer_cols()
            htot = nRows * answer_height + (nRows-1) * answer_inter_height
            (_, ay, _, ah) = self.st.get_answer_bbox()
            hoffset = (ah - htot) / 2
            
            # Add propositions
            for j in range(nQ):
                # Define col
                id_col = j % 2
                id_row = j // 2
                
                # Fix width
                if id_col == 0:
                    left = col1
                else:
                    left = col2
            
                # Add background
                local_offset = id_row * (answer_height + answer_inter_height)
                shapes =  slide.shapes
                
                # Check if one or more answers
                a = answers[j]
    
                if answers_id is not None and j in answers_id:
                    color_bg = RGBColor.from_string(self.st.COLOR_BBOX_BACKGROUND_CORRECT)
                else:
                    color_bg = RGBColor.from_string(self.st.COLOR_BBOX_BACKGROUND)
                    
                self._add_frame(
                    shapes=shapes,
                    bbox=[                
                        Mm(left), 
                        Mm(local_offset + ay + hoffset), 
                        Mm(col), 
                        Mm(answer_height)
                    ],
                    text = None,
                    color_bg=color_bg,
                    color_line=color_bg,
                    font_file=pfont,
                )
                            
                self._add_frame(
                    shapes=shapes,
                    bbox=[                
                        Mm(left + im), 
                        Mm(local_offset + ay + hoffset + im), 
                        Mm(col - 2*im), 
                        Mm(answer_height - 2*im)                        
                    ],
                    text=a,
                    color_bg=None,
                    color_line=None,
                    font=PptxFactory.FONT_ANSWER,
                    alignement=PP_ALIGN.LEFT,
                    font_file=pfont,
                )
                                
            