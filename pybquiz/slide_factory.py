from pptx.util import Mm
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.presentation import Presentation as PObject
import numpy as np


def pos_to_char(pos):
    return chr(pos + 97)

class SlideFactory:
    
    TITLE_SUBTITLE = 0
    TITLE_CONTENT = 1
    SECTION_HEADER = 2
    TWO_CONTENT = 3
    COMPARISON = 4
    TITLE = 5
    BLANK = 6
    CONTENT_CAPTION = 7
    IMAGE_CAPTION = 8
        
    # Base slide frame in mm
    S_FRAME_WITDH = 254  # 10 Inches
    S_FRAME_HEIGH = 190.5  # 7.5 Inches
    S_FRAME_MARGIN = 20
    S_FRAME_MARGIN_TINY = 5
    # Question frame
    S_QUEST_BBOX_W = S_FRAME_WITDH - (2 * S_FRAME_MARGIN)
    S_QUEST_BBOX_H = 30
    # Answer frames
    S_ANSWERS_WIDTH = S_QUEST_BBOX_W
    S_ANSWERS_TOP = S_QUEST_BBOX_H + 2*S_FRAME_MARGIN
    S_ANSWERS_HEIGHT = S_FRAME_HEIGH - S_ANSWERS_TOP - S_FRAME_MARGIN
    
    S_ANSWER_COL_MARGIN = S_FRAME_MARGIN / 2
    S_ANSWER_COL1_LEFT = S_FRAME_MARGIN
    S_ANSWER_COL_WIDTH = (S_ANSWERS_WIDTH - S_ANSWER_COL_MARGIN) / 2
    S_ANSWER_COL2_LEFT = S_FRAME_MARGIN + S_ANSWER_COL_WIDTH + S_ANSWER_COL_MARGIN
    S_ANSWER_ROW_HEIGHT_MAX = 20
    S_ANSWER_ROW_HEIGHT_MARGIN = 10

    # Font
    FONT_TITLE = 40
    FONT_QUESTION = 28
    FONT_ANSWER = 20
    
    # COLORS
    COLOR_ANSWER_CORRECT = RGBColor.from_string("8fce00")
    COLOR_BBOX_BACKGROUND = RGBColor.from_string("536878")
    COLOR_BBOX_LINE = RGBColor.from_string("36454F")
    COLOR_TEXT = RGBColor.from_string("EBEBEB")
    
    @staticmethod
    def get_system_fonts():
        return "/usr/share/fonts/opentype/fonts-hosny-amiri/Amiri-Regular.ttf"
    
    @staticmethod
    def add_title_subtitle(root: PObject, title: str = "", subtitle: str = "", img_bg: str = None):
        
            # Add slide
            title_slide = root.slides.add_slide(root.slide_layouts[SlideFactory.BLANK]) 
            shapes =  title_slide.shapes            
                        
            # Add background
            if img_bg is not None:
                shapes.add_picture(img_bg, Mm(0), Mm(0), height=Mm(SlideFactory.S_FRAME_HEIGH))
            
            shape = shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, 
                Mm(SlideFactory.S_FRAME_MARGIN), 
                Mm(3*SlideFactory.S_FRAME_MARGIN), 
                Mm(SlideFactory.S_FRAME_WITDH-2*SlideFactory.S_FRAME_MARGIN), 
                Mm(2*SlideFactory.S_FRAME_MARGIN),
            )
            
            shape.fill.solid()
            shape.fill.fore_color.rgb = SlideFactory.COLOR_BBOX_BACKGROUND
            shape.line.fill.solid()
            shape.line.fill.fore_color.rgb = SlideFactory.COLOR_BBOX_LINE
            shape.shadow.inherit = False
                        
            tf = shape.text_frame
            tf.text = title
            tf.fit_text(max_size=SlideFactory.FONT_TITLE, font_file=SlideFactory.get_system_fonts()) 
            tf.paragraphs[0].font.color.rgb = SlideFactory.COLOR_TEXT
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # shapeOut.fill.solid()
            # shapeOut.fill.fore_color.rgb = SlideFactory.COLOR_BBOX_BACKGROUND
            # shapeOut.line.fill.solid()
            # shapeOut.line.fill.fore_color.rgb = SlideFactory.COLOR_BBOX_LINE
            # shapeOut.shadow.inherit = False
            
            # First place holder (title)
            # Second place holder (subtitle)
            # title_slide.placeholders[0].text = title
            
            # title_slide.placeholders[1].text = subtitle       

            
    @staticmethod
    def add_question(root: PObject, i: int = 1, question: str = "", answers: list[str] = [], answers_id: list[int] = None, img_bg: str = None):
        
            # Add slide and question            
            question_slide = root.slides.add_slide(root.slide_layouts[SlideFactory.BLANK]) 
            
            # Add text shape
            shapes =  question_slide.shapes

            if img_bg is not None:
                shapes.add_picture(img_bg, Mm(0), Mm(0), height=Mm(SlideFactory.S_FRAME_HEIGH))
                
            shapeOut = shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, 
                Mm(SlideFactory.S_FRAME_MARGIN), 
                Mm(SlideFactory.S_FRAME_MARGIN), 
                Mm(SlideFactory.S_QUEST_BBOX_W), 
                Mm(SlideFactory.S_QUEST_BBOX_H),
            )
            
            shapeOut.fill.solid()
            shapeOut.fill.fore_color.rgb = SlideFactory.COLOR_BBOX_BACKGROUND
            shapeOut.line.fill.solid()
            shapeOut.line.fill.fore_color.rgb = SlideFactory.COLOR_BBOX_LINE
            shapeOut.shadow.inherit = False
            
            shape = shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, 
                Mm(SlideFactory.S_FRAME_MARGIN + SlideFactory.S_FRAME_MARGIN_TINY), 
                Mm(SlideFactory.S_FRAME_MARGIN + SlideFactory.S_FRAME_MARGIN_TINY), 
                Mm(SlideFactory.S_QUEST_BBOX_W - 2*SlideFactory.S_FRAME_MARGIN_TINY), 
                Mm(SlideFactory.S_QUEST_BBOX_H - 2*SlideFactory.S_FRAME_MARGIN_TINY),
            )
            shape.shadow.inherit = False                
            shape.fill.background()
            shape.line.fill.background()
            
            tf = shape.text_frame
            tf.text = "Q{}: {}".format(i, question)
            tf.fit_text(max_size=SlideFactory.FONT_QUESTION, font_file=SlideFactory.get_system_fonts()) 
            tf.paragraphs[0].font.color.rgb = SlideFactory.COLOR_TEXT

            SlideFactory._add_answers(question_slide, answers, answers_id)
            
            
    @staticmethod
    def _add_answers(slide, answers: list[str], answers_id: list[int] = None):
                           
            nQ = len(answers)
            nRows = np.ceil(nQ/2).astype(int)
            
            # # Add background
            # shapes =  slide.shapes
            # shape = shapes.add_shape(
            #     MSO_SHAPE.ROUNDED_RECTANGLE, 
            #     Mm(SlideFactory.S_FRAME_MARGIN), 
            #     Mm(SlideFactory.S_ANSWERS_TOP), 
            #     Mm(SlideFactory.S_ANSWERS_WIDTH), 
            #     Mm(SlideFactory.S_ANSWERS_HEIGHT),
            # )
            # shape.fill.solid()
            # shape.fill.fore_color.rgb = SlideFactory.COLOR_BBOX_DEBUG
                        
            # Check if only one proposition and no answer
            if nQ <= 1 and answers_id is None:
                return
            
            # Get total heigh based on the number of entries
            htot = nRows * SlideFactory.S_ANSWER_ROW_HEIGHT_MAX + (nRows-1) * SlideFactory.S_ANSWER_ROW_HEIGHT_MARGIN
            hoffset = (SlideFactory.S_ANSWERS_HEIGHT - htot) / 2
            
            # Add propositions
            for j in range(nQ):
                # Define col
                id_col = j % 2
                id_row = j // 2
                
                # Fix width
                if id_col == 0:
                    left = SlideFactory.S_ANSWER_COL1_LEFT
                else:
                    left = SlideFactory.S_ANSWER_COL2_LEFT
            
                # Add background
                local_offset = id_row * (SlideFactory.S_ANSWER_ROW_HEIGHT_MAX + SlideFactory.S_ANSWER_ROW_HEIGHT_MARGIN)
                shapes =  slide.shapes
                
                # Add background
                shapes =  slide.shapes
                
                shapeOut = shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE, 
                    Mm(left), 
                    Mm(SlideFactory.S_ANSWERS_TOP + local_offset + hoffset), 
                    Mm(SlideFactory.S_ANSWER_COL_WIDTH), 
                    Mm(SlideFactory.S_ANSWER_ROW_HEIGHT_MAX)
                )
                shapeOut.fill.solid()
                shapeOut.fill.fore_color.rgb = SlideFactory.COLOR_BBOX_BACKGROUND
                shapeOut.line.fill.solid()
                shapeOut.line.fill.fore_color.rgb = SlideFactory.COLOR_BBOX_LINE
                
                shape = shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE, 
                    Mm(left), 
                    Mm(SlideFactory.S_ANSWERS_TOP + local_offset + hoffset), 
                    Mm(SlideFactory.S_ANSWER_COL_WIDTH), 
                    Mm(SlideFactory.S_ANSWER_ROW_HEIGHT_MAX)
                )
                shape.fill.background()
                shape.line.fill.background()                                            
                tf = shape.text_frame

                # Check if one or more answers
                a = answers[j]
                if nQ > 1:
                    a = "{}) {}".format(pos_to_char(j), a)
                # Set text
                tf.text = a
                tf.fit_text(max_size=SlideFactory.FONT_ANSWER, font_file=SlideFactory.get_system_fonts()) 
                tf.paragraphs[0].font.color.rgb = SlideFactory.COLOR_TEXT
                                
                if answers_id is not None and j in answers_id:
                    tf.paragraphs[0].font.color.rgb = SlideFactory.COLOR_ANSWER_CORRECT
            
            